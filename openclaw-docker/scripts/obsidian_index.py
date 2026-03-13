#!/usr/bin/env python3
"""
Obsidian FTS5 Indexer — fast text extraction, no ML needed for search indexing.

Formats:
  .md   → native read
  .pdf  → pypdfium2  (installed with docling, ~1s per 1000 pages)
  .docx → python-docx (installed with docling)
  .xlsx → openpyxl   (installed with docling)
  .pptx → python-pptx (installed with docling)

Docling (ML, slow) is NOT used here — it's for document_intelligence skill
when the bot needs to read a specific complex document with full structure.
"""

import os
import sys
import sqlite3
import re
import time
import argparse
from pathlib import Path
from tqdm import tqdm

# System and User Vault Paths
sys_vault_env = os.environ.get("SYSTEM_VAULT_PATH", "/data/obsidian/vault")
user_vault_env = os.environ.get("USER_VAULT_PATH", "/data/abror_vault")

VAULT_PATHS = [Path(sys_vault_env)]
# Only add user vault if it exists, to avoid errors if not mapped
if Path(user_vault_env).exists() and user_vault_env != sys_vault_env:
    VAULT_PATHS.append(Path(user_vault_env))

DB_DEFAULT = Path(sys_vault_env) / "Bot" / "obsidian.db"
MAX_CHUNK_LINES = 200

DB_PATH = DB_DEFAULT


# ── Extractors ────────────────────────────────────────────────────────────────

def extract_pdf(path: Path) -> str:
    """Extract PDF to Markdown using MarkItDown (better table support)"""
    from markitdown import MarkItDown
    md = MarkItDown(enable_plugins=False)
    result = md.convert(str(path))
    return result.text_content if result.text_content else ""


def extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(str(path))
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            # Detect headings by style
            if para.style.name.startswith("Heading"):
                level = para.style.name.split()[-1] if para.style.name.split()[-1].isdigit() else "1"
                parts.append(f"{'#' * int(level)} {para.text}")
            else:
                parts.append(para.text)
    # Tables
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(c.text.strip() for c in row.cells if c.text.strip()))
    return "\n\n".join(parts)


def extract_xlsx(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
    parts = []
    for sheet in wb.sheetnames:
        ws = wb[sheet]
        parts.append(f"## {sheet}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) for c in row if c is not None]
            if cells:
                parts.append(" | ".join(cells))
    wb.close()
    return "\n\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"## Slide {i}")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                parts.append(shape.text)
    return "\n\n".join(parts)


EXTRACTORS = {
    ".pdf":  extract_pdf,
    ".docx": extract_docx,
    ".xlsx": extract_xlsx,
    ".pptx": extract_pptx,
}


# ── DB ────────────────────────────────────────────────────────────────────────

def get_db():
    try:
        DB_PATH.parent.mkdir(parents=True, exist_ok=True)
        conn = sqlite3.connect(DB_PATH)
        conn.execute("PRAGMA journal_mode=WAL")
        conn.execute("CREATE TABLE IF NOT EXISTS files (path TEXT PRIMARY KEY, mtime REAL)")
        conn.execute("""
            CREATE VIRTUAL TABLE IF NOT EXISTS chunks USING fts5(
                file_path, section_title, content,
                tokenize='unicode61 remove_diacritics 1'
            )
        """)
        conn.commit()
        return conn
    except Exception as e:
        print(f"CRITICAL: Database error: {e}")
        sys.exit(1)


# ── Chunking ──────────────────────────────────────────────────────────────────

def chunk_markdown(text, file_path):
    lines = text.splitlines()
    chunks = []
    current_title = "(intro)"
    current_lines = []

    for line in lines:
        header = re.match(r'^(#{1,4})\s+(.+)', line)
        if header:
            if current_lines:
                chunks.append((file_path, current_title, "\n".join(current_lines)))
            current_title = header.group(2).strip()
            current_lines = [line]
        else:
            current_lines.append(line)
            if len(current_lines) >= MAX_CHUNK_LINES:
                chunks.append((file_path, current_title + " [cont]", "\n".join(current_lines)))
                current_lines = []

    if current_lines:
        chunks.append((file_path, current_title, "\n".join(current_lines)))
    return chunks


# ── Index ─────────────────────────────────────────────────────────────────────

def index(vault_path=None, force=False):
    vaults = [Path(vault_path)] if vault_path else VAULT_PATHS
    print(f"\n📂 Vaults: {[str(v) for v in vaults]}")

    conn = get_db()
    cur = conn.cursor()

    if force:
        cur.execute("DELETE FROM chunks")
        cur.execute("DELETE FROM files")
        conn.commit()
        print("🗑  Force mode: index cleared")

    extensions = list(EXTRACTORS.keys()) + ["*.md"]
    all_files = []
    for vault in vaults:
        for ext in extensions:
            all_files.extend(vault.rglob(f"*{ext}" if not ext.startswith("*") else ext))

    to_index, to_skip = [], 0
    for f in all_files:
        rel = str(f)  # store absolute paths inside container
        mtime = f.stat().st_mtime
        row = cur.execute("SELECT mtime FROM files WHERE path=?", (rel,)).fetchone()
        if not force and row and abs(row[0] - mtime) < 1:
            to_skip += 1
        else:
            to_index.append((f, rel))

    print(f"📊 {len(all_files)} files total: {len(to_index)} to index, {to_skip} unchanged")

    if not to_index:
        print("✅ Everything up to date!")
        conn.close()
        return

    added = updated = errors = 0
    t_total = time.time()

    with tqdm(to_index, unit="file", ncols=90, colour="green") as pbar:
        for doc_file, rel in pbar:
            ext = doc_file.suffix.lower()
            size_mb = doc_file.stat().st_size / 1_048_576
            pbar.set_description(f"{ext[1:].upper():5s}")
            pbar.set_postfix_str(f"{doc_file.name[:40]} ({size_mb:.1f}MB)", refresh=True)

            t0 = time.time()
            try:
                if ext == ".md":
                    text = doc_file.read_text(encoding="utf-8", errors="ignore")
                elif ext in EXTRACTORS:
                    text = EXTRACTORS[ext](doc_file)
                else:
                    continue
            except Exception as e:
                tqdm.write(f"  ❌ {doc_file.name}: {e}")
                errors += 1
                continue

            if not text or not text.strip():
                errors += 1
                continue

            elapsed = time.time() - t0
            pbar.set_postfix_str(
                f"{doc_file.name[:30]} ✅ {len(text):,}ch {elapsed:.1f}s", refresh=True
            )

            row = cur.execute("SELECT mtime FROM files WHERE path=?", (rel,)).fetchone()
            mtime = doc_file.stat().st_mtime
            cur.execute("DELETE FROM chunks WHERE file_path=?", (rel,))
            chunks = chunk_markdown(text, rel)
            cur.executemany(
                "INSERT INTO chunks(file_path, section_title, content) VALUES(?,?,?)", chunks
            )
            if row:
                cur.execute("UPDATE files SET mtime=? WHERE path=?", (mtime, rel))
                updated += 1
            else:
                cur.execute("INSERT INTO files(path, mtime) VALUES(?,?)", (rel, mtime))
                added += 1
            conn.commit()

    elapsed_total = time.time() - t_total
    conn.close()
    print(f"\n✅ Done in {elapsed_total:.1f}s — +{added} new, ~{updated} updated, {errors} errors, {to_skip} skipped")


# ── Search ────────────────────────────────────────────────────────────────────

def search(query, limit=3):
    if not DB_PATH.exists():
        print("Index not found. Run without --search to build it.")
        sys.exit(1)

    conn = sqlite3.connect(DB_PATH)
    conn.row_factory = sqlite3.Row
    try:
        rows = conn.execute("""
            SELECT file_path, section_title, content, rank
            FROM chunks WHERE chunks MATCH ? ORDER BY rank LIMIT ?
        """, (query, limit)).fetchall()
    except sqlite3.OperationalError:
        rows = conn.execute("""
            SELECT file_path, section_title, content, 0 as rank
            FROM chunks WHERE content LIKE ? LIMIT ?
        """, (f"%{query}%", limit)).fetchall()
    conn.close()

    if not rows:
        print(f"No results for: {query}")
        return

    print(f'\n🔍 Results for: "{query}"\n')
    for row in rows:
        print(f"📄 {row['file_path']}  [§ {row['section_title']}]")
        print("────────────────────────────────")
        print("\n".join(row['content'].splitlines()[:40]))
        print()


if __name__ == "__main__":
    parser = argparse.ArgumentParser()
    parser.add_argument("--search", "-s")
    parser.add_argument("--limit", "-l", type=int, default=3)
    parser.add_argument("--vault")
    parser.add_argument("--db")
    parser.add_argument("--force", "-f", action="store_true")
    args = parser.parse_args()

    if args.db:
        DB_PATH = Path(args.db)
    if args.search:
        search(args.search, args.limit)
    else:
        index(args.vault, force=args.force)
